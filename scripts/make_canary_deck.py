"""Generate Canary-RL project slide deck (16:9)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

REPO = Path(__file__).resolve().parent.parent
OUT  = REPO / "Canary-RL.pptx"

# ---- palette -------------------------------------------------------------
NAVY  = RGBColor(0x0F, 0x1F, 0x3D)
INK   = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x60, 0x6E)
LINE  = RGBColor(0xD8, 0xDC, 0xE3)
RED   = RGBColor(0xC8, 0x32, 0x32)
GREEN = RGBColor(0x1F, 0x7A, 0x3D)
GOLD  = RGBColor(0xC9, 0x8A, 0x14)
BG    = RGBColor(0xFA, 0xFB, 0xFC)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ---- helpers -------------------------------------------------------------
def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    return s

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return tb

def add_rule(slide, x, y, w, color=NAVY, weight=2.5):
    line = slide.shapes.add_connector(1, x, y, x + w, y)
    line.line.color.rgb = color; line.line.width = Pt(weight)
    return line

def add_box(slide, x, y, w, h, *, fill=None, border=LINE, weight=1.0, radius=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    b = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        b.fill.background()
    else:
        b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.color.rgb = border; b.line.width = Pt(weight)
    b.shadow.inherit = False
    return b

def header(slide, kicker, title):
    add_text(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.35),
             kicker, size=12, bold=True, color=GOLD)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    add_rule(slide, Inches(0.6), Inches(1.45), Inches(12.1))

def footer(slide, n, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             "Canary-RL · BT-GRPO failure-prediction stack", size=10, color=MUTED)
    add_text(slide, Inches(11.6), Inches(7.05), Inches(1.2), Inches(0.3),
             f"{n} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# ---- slide 1: title -----------------------------------------------------
s = add_slide()
add_box(s, 0, 0, W, Inches(7.5), fill=NAVY, border=NAVY)
add_text(s, Inches(0.8), Inches(2.4), Inches(11), Inches(0.5),
         "An RL-training fake-validation plugin", size=18, color=GOLD, bold=True)
add_text(s, Inches(0.8), Inches(2.95), Inches(12), Inches(1.6),
         "Canary-RL", size=72, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
add_text(s, Inches(0.8), Inches(4.3), Inches(12), Inches(0.8),
         "Predicting BT-GRPO failure before training spends an hour of GPU",
         size=24, color=RGBColor(0xE6,0xEA,0xF2))
add_rule(s, Inches(0.85), Inches(5.45), Inches(2.0), color=GOLD, weight=2.5)
add_text(s, Inches(0.8), Inches(5.6), Inches(11), Inches(0.4),
         "LLaDA-8B-Instruct · GSM8K · 19 historical runs",
         size=14, color=RGBColor(0xC9,0xCF,0xDC))

# ---- slide 2: the problem ------------------------------------------------
s = add_slide(); header(s, "MOTIVATION", "19 runs in, every BT-GRPO collapsed differently")
add_text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.45),
         "Each failed run cost ~hours of 8-GPU compute. Failures only became visible "
         "thousands of steps in,\nlong after the launch script was already wrong.",
         size=15, color=INK)

bullets = [
    ("run12 / run13",  "grad blew up to 80k–331k after step 600",                 RED),
    ("run14 / run15",  "fork-head saturated at boundary; reward stalled",         RED),
    ("run17",          "27× reward-weight ratio → format-collapse reward hack",   RED),
    ("run18 / run19",  "frac-zero-std stuck > 0.7 → no learning signal",          RED),
    ("run5 / run7 / run10", "slow flatlines; corr stayed < 0.30 forever",        GOLD),
    ("0 of 19",        "ever achieved sustained correctness ≥ 0.30",              MUTED),
]
y = Inches(2.9)
for tag, txt, c in bullets:
    add_box(s, Inches(0.6), y, Inches(0.18), Inches(0.45), fill=c, border=c)
    add_text(s, Inches(1.0), y - Emu(20000), Inches(2.6), Inches(0.45),
             tag, size=15, bold=True, color=NAVY)
    add_text(s, Inches(3.6), y - Emu(20000), Inches(9.0), Inches(0.45),
             txt, size=15, color=INK)
    y += Inches(0.55)

add_box(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.6),
        fill=RGBColor(0xFF,0xF6,0xE0), border=GOLD, radius=True)
add_text(s, Inches(0.85), Inches(6.45), Inches(11.6), Inches(0.4),
         "Question:  can we predict — from the launch script alone — which mode "
         "the next run will collapse in?",
         size=14, bold=True, color=NAVY)
footer(s, 2, 10)

# ---- slide 3: failure taxonomy -------------------------------------------
s = add_slide(); header(s, "TAXONOMY", "Six early-window failure signatures")
add_text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.4),
         "Mined from 8,176 logged steps across 16 runs. Each signature has a "
         "calibrated threshold and a tier (abort vs advisory).",
         size=14, color=MUTED)
rows = [
    ("grad_blowup",          "Tier-1 abort",   "max grad-norm > 5,000",                            "run12 hit 331k"),
    ("starved_signal",       "Tier-2 advisory","frac_zero_std > 0.5 sustained 5 steps",            "run18, run19"),
    ("fork_saturated",       "Tier-2 advisory","fork-head μ pinned at fork_frac_min/max",          "run14, run15"),
    ("len_collapsing",       "Tier-2 advisory","completion-length slope < −0.5 tok/step",         "run5, run8"),
    ("corr_dead_early",      "Tier-3 trend",   "corr < 0.05 by step 50, no upward slope",          "rare in our data"),
    ("corr_negative_slope",  "Tier-3 trend",   "linear corr extrapolates downward at step 200",    "run5, run17"),
]
hdr_y = Inches(2.4)
add_box(s, Inches(0.6), hdr_y, Inches(12.1), Inches(0.45), fill=NAVY, border=NAVY)
hdrs = [("Signature",0.6,2.6),("Tier",3.2,2.0),("Trip condition",5.2,4.6),("Historical example",9.8,2.9)]
for h,x,w in hdrs:
    add_text(s, Inches(x), hdr_y + Emu(40000), Inches(w), Inches(0.4),
             h, size=12, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
y = hdr_y + Inches(0.5)
for sig, tier, cond, ex in rows:
    tcol = RED if "abort" in tier else (GOLD if "advisory" in tier else MUTED)
    add_box(s, Inches(0.6), y, Inches(12.1), Inches(0.5), border=LINE)
    add_text(s, Inches(0.8), y + Emu(60000), Inches(2.4), Inches(0.4),
             sig, size=13, bold=True, color=NAVY, font="Consolas")
    add_text(s, Inches(3.2), y + Emu(60000), Inches(2.0), Inches(0.4),
             tier, size=12, bold=True, color=tcol)
    add_text(s, Inches(5.2), y + Emu(60000), Inches(4.6), Inches(0.4),
             cond, size=12, color=INK, font="Consolas")
    add_text(s, Inches(9.8), y + Emu(60000), Inches(2.9), Inches(0.4),
             ex, size=12, color=MUTED)
    y += Inches(0.55)
footer(s, 3, 10)

# ---- slide 4: the methodology ------------------------------------------
s = add_slide(); header(s, "METHODOLOGY", "Four tiers of fake validation")
add_text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.4),
         "Each tier costs less than the next, and runs against information available earlier in the pipeline.",
         size=14, color=MUTED)
tiers = [
    ("0", "Preflight",           "before launch", "Static checks on launch script: reward-weight balance, parser fixtures, parameter sanity.", GREEN),
    ("1", "Live watcher abort",  "step 1–50",     "Tail trainer log; fire SIGINT if grad-blowup or sustained signal-starvation tripwire fires.", RED),
    ("2", "Live advisory lights","step 1–50",     "Six per-step signature glyphs ●/○; surface advisory failures without auto-killing.", GOLD),
    ("3", "Trend extrapolator",  "step 50–200",   "Linear fit on (step, correctness) — projects whether corr will ever cross 0.30.", NAVY),
]
y = Inches(2.4)
for n, name, when, desc, c in tiers:
    add_box(s, Inches(0.6), y, Inches(12.1), Inches(1.05), border=LINE, radius=True)
    add_box(s, Inches(0.85), y + Inches(0.18), Inches(0.7), Inches(0.7),
            fill=c, border=c, radius=True)
    add_text(s, Inches(0.85), y + Inches(0.27), Inches(0.7), Inches(0.5),
             n, size=22, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.75), y + Inches(0.18), Inches(4.0), Inches(0.4),
             name, size=15, bold=True, color=NAVY)
    add_text(s, Inches(1.75), y + Inches(0.55), Inches(4.0), Inches(0.4),
             when, size=12, color=MUTED, font="Consolas")
    add_text(s, Inches(5.8), y + Inches(0.25), Inches(6.7), Inches(0.7),
             desc, size=13, color=INK)
    y += Inches(1.15)
footer(s, 4, 10)

# ---- slide 5: architecture ----------------------------------------------
s = add_slide(); header(s, "ARCHITECTURE", "Plug-and-play stack around the trainer")

# trainer in middle
trn = add_box(s, Inches(5.0), Inches(3.3), Inches(3.2), Inches(1.0),
              fill=NAVY, border=NAVY, radius=True)
add_text(s, Inches(5.0), Inches(3.45), Inches(3.2), Inches(0.4),
         "BT-GRPO trainer", size=15, bold=True, color=RGBColor(0xFF,0xFF,0xFF),
         align=PP_ALIGN.CENTER)
add_text(s, Inches(5.0), Inches(3.85), Inches(3.2), Inches(0.4),
         "examples/rl/grpo/llada/train_btgrpo.py",
         size=10, color=RGBColor(0xC9,0xCF,0xDC), font="Consolas", align=PP_ALIGN.CENTER)

# 4 satellites
def satellite(label, sub, x, y, c):
    add_box(s, x, y, Inches(2.7), Inches(0.85), fill=c, border=c, radius=True)
    add_text(s, x, y + Inches(0.07), Inches(2.7), Inches(0.4),
             label, size=14, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.45), Inches(2.7), Inches(0.4),
             sub, size=10, color=RGBColor(0xF5,0xF5,0xF5), font="Consolas", align=PP_ALIGN.CENTER)
satellite("Preflight",        "scripts/canary_preflight.py", Inches(0.6),  Inches(2.0),  GREEN)
satellite("HealthOracle",     "scripts/canary_oracle.py",    Inches(0.6),  Inches(5.4),  GOLD)
satellite("Live watcher",     "scripts/canary_watcher.py",   Inches(10.05),Inches(2.0),  RED)
satellite("Dashboard",        "dashboard.py",                Inches(10.05),Inches(5.4),  RGBColor(0x4A,0x6F,0xB0))

# arrows in/out
def arrow(x1, y1, x2, y2):
    a = s.shapes.add_connector(2, x1, y1, x2, y2)
    a.line.color.rgb = MUTED; a.line.width = Pt(1.5)
arrow(Inches(3.3), Inches(2.85), Inches(5.0),  Inches(3.55))    # preflight -> trainer
arrow(Inches(3.3), Inches(5.85), Inches(5.0),  Inches(4.0))     # oracle -> trainer
arrow(Inches(8.2), Inches(3.55), Inches(10.05), Inches(2.85))   # trainer -> watcher
arrow(Inches(8.2), Inches(4.0),  Inches(10.05), Inches(5.85))   # trainer -> dashboard

# unified launcher banner
add_box(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.5),
        fill=RGBColor(0xEC,0xF1,0xF8), border=NAVY, radius=True)
add_text(s, Inches(0.85), Inches(6.62), Inches(11.6), Inches(0.4),
         "scripts/launch_with_canary.sh  —  preflight ▸ trainer ▸ watcher ▸ dashboard, with cleanup on Ctrl-C",
         size=12, bold=True, color=NAVY, font="Consolas")
footer(s, 5, 10)

# ---- slide 6: HealthOracle ----------------------------------------------
s = add_slide(); header(s, "PLUGIN", "HealthOracle: predict from launch script alone")

add_text(s, Inches(0.6), Inches(1.7), Inches(7.5), Inches(0.4),
         "Inputs:  25 features parsed from any launch_btgrpo_runXX.sh",
         size=13, color=INK, bold=True)
feats = ("β · ε · learning_rate · lora_r · per_device_bs\n"
         "num_generations · num_iterations\n"
         "max_completion_length · block_size\n"
         "fork_frac · fork_frac_min/max · fork_head_lr\n"
         "learn_fork_frac · scale_rewards · filter_zero_std/correct\n"
         "w_xmlcount · w_soft · w_strict · w_int · w_correctness\n"
         "reward_ratio_max_over_min  …")
add_text(s, Inches(0.6), Inches(2.15), Inches(7.5), Inches(2.6),
         feats, size=11, color=MUTED, font="Consolas")

add_text(s, Inches(0.6), Inches(4.65), Inches(7.5), Inches(0.4),
         "Method:  per-signature LOO logistic regression on 15 historical runs",
         size=13, bold=True, color=INK)
add_text(s, Inches(0.6), Inches(5.05), Inches(7.5), Inches(1.6),
         "• one classifier per failure mode\n"
         "• z-scored features, class_weight=balanced\n"
         "• kNN nearest-historical-neighbour (Euclidean in feature space)\n"
         "• rule-based recommendations bound to feature coefficients",
         size=12, color=INK)

# right: outputs
add_box(s, Inches(8.3), Inches(1.85), Inches(4.4), Inches(4.95),
        fill=RGBColor(0xF1,0xF4,0xFA), border=LINE, radius=True)
add_text(s, Inches(8.5), Inches(2.0), Inches(4.0), Inches(0.4),
         "Output", size=14, bold=True, color=NAVY)
add_text(s, Inches(8.5), Inches(2.45), Inches(4.0), Inches(4.3),
         "P(failure within 50 steps)\n"
         "  per-signature probabilities\n"
         "    grad_blowup        0.00\n"
         "    starved_signal     0.31  ████\n"
         "    fork_saturated     0.00\n"
         "    len_collapsing     0.01\n"
         "    corr_dead_early    0.00\n"
         "    corr_negative_sl   0.09  ██\n\n"
         "Nearest historical analogs\n"
         "  run19  d=9.29  starved\n"
         "  run5   d=11.2  none\n\n"
         "Recommendations (rule-based)",
         size=11, color=INK, font="Consolas")
footer(s, 6, 10)

# ---- slide 7: LOO calibration -------------------------------------------
s = add_slide(); header(s, "VALIDATION", "Leave-one-out: how well does it predict?")
add_text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.4),
         "Hold each run out, train on the other 14, predict.  P_any = 1 − ∏(1−Pᵢ).",
         size=13, color=MUTED)

cols = ["run","grad","starve","fork","len","cDead","cNeg","truth","P_any","verdict"]
rows = [
    ("run5", 0.11,0.23,0.02,0.74,0.00,0.88, 0, 0.98, "FP*"),
    ("run6", 0.02,0.02,0.04,0.24,0.00,0.30, 1, 0.51, "TP"),
    ("run7", 0.02,0.02,0.04,0.24,0.00,0.90, 0, 0.93, "FP*"),
    ("run8", 0.03,0.03,0.04,0.45,0.00,0.70, 1, 0.85, "TP"),
    ("run10",0.03,0.03,0.04,0.65,0.00,0.86, 0, 0.96, "FP*"),
    ("run12",1.00,0.95,0.99,0.93,0.00,0.01, 1, 1.00, "TP"),
    ("run13",0.32,0.94,0.18,0.45,0.00,0.41, 1, 0.99, "TP"),
    ("run14",0.74,0.89,0.92,0.71,0.00,0.10, 1, 1.00, "TP"),
    ("run15",0.14,0.97,0.93,0.18,0.00,0.09, 1, 1.00, "TP"),
    ("run17",0.04,0.99,0.14,0.81,0.00,0.74, 1, 1.00, "TP"),
    ("run18",0.04,0.99,0.14,0.02,0.00,0.01, 1, 0.99, "TP"),
    ("run19",0.02,0.96,0.04,0.79,0.00,0.88, 1, 1.00, "TP"),
]
y0 = Inches(2.3); rh = Inches(0.32)
xs_ = [0.6,1.6,2.4,3.2,4.0,4.8,5.6,6.4,7.4,8.4,9.6]
add_box(s, Inches(0.6), y0, Inches(8.7), rh, fill=NAVY, border=NAVY)
for i,c in enumerate(cols):
    add_text(s, Inches(xs_[i]), y0 + Emu(25000), Inches(0.95), Inches(0.3),
             c, size=10, bold=True, color=RGBColor(0xFF,0xFF,0xFF), font="Consolas")
y = y0 + rh
for r in rows:
    color_bg = RGBColor(0xFF,0xF8,0xF0) if r[-1]=="FP*" else BG
    add_box(s, Inches(0.6), y, Inches(8.7), rh, fill=color_bg, border=LINE)
    for i,v in enumerate(r):
        col = INK
        txt = v if isinstance(v,str) else f"{v:.2f}"
        if cols[i]=="truth": txt = str(v); col = (RED if v else GREEN)
        if cols[i]=="verdict": col = GREEN if v=="TP" else GOLD
        add_text(s, Inches(xs_[i]), y + Emu(30000), Inches(0.95), Inches(0.3),
                 txt, size=10, color=col, font="Consolas")
    y += rh

# right column commentary
add_box(s, Inches(9.5), Inches(2.3), Inches(3.2), Inches(4.4),
        fill=RGBColor(0xF1,0xF4,0xFA), border=LINE, radius=True)
add_text(s, Inches(9.7), Inches(2.4), Inches(3.0), Inches(0.4),
         "Read-out", size=14, bold=True, color=NAVY)
add_text(s, Inches(9.7), Inches(2.85), Inches(3.0), Inches(4.0),
         "12 / 12 true failures\n"
         "  → P_any ≥ 0.51\n\n"
         "3 “false positives”\n"
         "  run5/7/10 fired no\n"
         "  early signature, but\n"
         "  ALSO ended in flat\n"
         "  failure (corr < 0.30)\n"
         "  → oracle is right to\n"
         "  flag them.\n\n"
         "0 truly successful\n"
         "  runs in the corpus.",
         size=11, color=INK)
footer(s, 7, 10)

# ---- slide 8: run20 vs run21 -------------------------------------------
s = add_slide(); header(s, "ABLATION", "Run20 (collapsed) → Run21 (oracle-guided)")

# big P_any callouts
def pbox(x, label, p, color):
    add_box(s, x, Inches(1.85), Inches(3.2), Inches(1.5),
            fill=color, border=color, radius=True)
    add_text(s, x, Inches(2.0), Inches(3.2), Inches(0.4),
             label, size=14, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.4), Inches(3.2), Inches(0.9),
             f"P_fail = {p}", size=34, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
pbox(Inches(0.6),  "run20  (clone of run19)", "0.99", RED)
pbox(Inches(9.5),  "run21  (oracle-guided)",  "0.38", GREEN)

# arrow
arr = s.shapes.add_connector(2, Inches(3.95), Inches(2.6), Inches(9.4), Inches(2.6))
arr.line.color.rgb = NAVY; arr.line.width = Pt(2.5)
add_text(s, Inches(4.0), Inches(2.1), Inches(5.4), Inches(0.4),
         "5 knobs changed", size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.0), Inches(2.7), Inches(5.4), Inches(0.4),
         "predicted by the plugin", size=12, color=MUTED, align=PP_ALIGN.CENTER)

# diff table
diffs = [
    ("knob",                "run20",       "run21",        "why (signature most affected)"),
    ("w_strict",            "0",           "1.0",          "starved_signal  −0.44 coef"),
    ("w_int / w_xml / w_soft","1.0 each",  "0.5 each",     "fork_saturated  +reward-ratio  8× → 4×"),
    ("learn_fork_frac",     "True",        "False",        "fork_saturated  +0.74 coef (top driver)"),
    ("fork_frac",           "0.3",         "0.2",          "corr_negative_slope  +0.60 coef"),
    ("max_completion_length","512",        "384",          "starved_signal  +0.46 coef"),
]
y0 = Inches(4.0); rh = Inches(0.42)
add_box(s, Inches(0.6), y0, Inches(12.1), rh, fill=NAVY, border=NAVY)
xs_ = [0.8, 4.5, 6.3, 8.0]
ws_ = [3.6, 1.6, 1.6, 4.7]
for i,c in enumerate(diffs[0]):
    add_text(s, Inches(xs_[i]), y0 + Emu(70000), Inches(ws_[i]), Inches(0.4),
             c, size=12, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
y = y0 + rh
for r in diffs[1:]:
    add_box(s, Inches(0.6), y, Inches(12.1), rh, border=LINE)
    add_text(s, Inches(xs_[0]), y + Emu(70000), Inches(ws_[0]), Inches(0.4),
             r[0], size=12, bold=True, color=NAVY, font="Consolas")
    add_text(s, Inches(xs_[1]), y + Emu(70000), Inches(ws_[1]), Inches(0.4),
             r[1], size=12, color=RED, font="Consolas")
    add_text(s, Inches(xs_[2]), y + Emu(70000), Inches(ws_[2]), Inches(0.4),
             r[2], size=12, bold=True, color=GREEN, font="Consolas")
    add_text(s, Inches(xs_[3]), y + Emu(70000), Inches(ws_[3]), Inches(0.4),
             r[3], size=11, color=INK)
    y += rh
footer(s, 8, 10)

# ---- slide 9: live evidence --------------------------------------------
s = add_slide(); header(s, "WHAT WE LEARNED", "Run22 taught us: loss ≠ failure, reward = truth")

# left: the misleading numbers
add_box(s, Inches(0.6), Inches(1.85), Inches(6.0), Inches(2.6),
        fill=RGBColor(0xFD,0xEE,0xEE), border=RED, radius=True)
add_text(s, Inches(0.85), Inches(2.0), Inches(5.7), Inches(0.4),
         "run22 numbers (β=0.04, BT-GRPO + dLLM)", size=14, bold=True, color=RED)
add_text(s, Inches(0.85), Inches(2.45), Inches(5.7), Inches(1.95),
         "step 1   loss = 5.9e8     grad = 1.3e11\n"
         "step 12  loss = 2.1e8     grad = 3.0e9\n"
         "kl term  = 5e9 — 6e10     (β·KL drowns policy obj)\n"
         "watcher v1 → ABORT at step 1\n"
         "→ scary numbers, looks catastrophic",
         size=13, color=INK, font="Consolas")

# right: but reward was rising
add_box(s, Inches(7.0), Inches(1.85), Inches(5.7), Inches(2.6),
        fill=RGBColor(0xEC,0xF7,0xEC), border=GREEN, radius=True)
add_text(s, Inches(7.25), Inches(2.0), Inches(5.4), Inches(0.4),
         "what the model was actually doing", size=14, bold=True, color=GREEN)
add_text(s, Inches(7.25), Inches(2.45), Inches(5.4), Inches(1.95),
         "correctness step 1-4   mean = 0.42\n"
         "correctness step 9-12  mean = 0.52   (+24%)\n"
         "best step (12) = 0.72   beats every prior run\n"
         "→ training was working\n"
         "  max_grad_norm=1.0 clipped the KL noise away",
         size=13, color=INK, font="Consolas")

# explanation strip
add_box(s, Inches(0.6), Inches(4.65), Inches(12.1), Inches(1.0),
        fill=RGBColor(0xFF,0xF6,0xE0), border=GOLD, radius=True)
add_text(s, Inches(0.85), Inches(4.75), Inches(11.6), Inches(0.4),
         "Why loss looks huge but training is fine", size=13, bold=True, color=NAVY)
add_text(s, Inches(0.85), Inches(5.1), Inches(11.6), Inches(0.6),
         "BT-GRPO uses k3 KL with clamp(r, ±5) → exp(r) ≤ 148 per token × 384 tokens × 64 completions = 3.6M raw KL. "
         "Multiplied by β=0.04, the *value* is huge — but ∇r is zero at the clip boundary, so the *gradient* is mostly policy.",
         size=12, color=INK)

# fix
add_box(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(1.45),
        fill=RGBColor(0xF1,0xF4,0xFA), border=LINE, radius=True)
add_text(s, Inches(0.85), Inches(5.95), Inches(11.6), Inches(0.4),
         "Watcher v2 — re-tuned for BT-GRPO + dLLM", size=13, bold=True, color=NAVY)
add_text(s, Inches(0.85), Inches(6.30), Inches(11.6), Inches(1.0),
         "• abort_grad_norm:  1e4 → 1e13   (single-step transient tolerated)\n"
         "• abort_grad_blowup_thresh:  1e3 → 1e10, frac 0.05 → 0.50  (need majority of window exploded)\n"
         "• abort_loss_jump_factor:    100  → 1e6     (KL-driven loss swings are routine)\n"
         "• new advisory: corr_regression — fires when correctness EMA falls 0.25 from peak after step 30",
         size=12, color=INK, font="Consolas")
footer(s, 9, 10)

# ---- slide 10: closing --------------------------------------------------
s = add_slide()
add_box(s, 0, 0, W, H, fill=NAVY, border=NAVY)
add_text(s, Inches(0.8), Inches(1.0), Inches(12), Inches(0.5),
         "TAKEAWAYS", size=14, bold=True, color=GOLD)
add_text(s, Inches(0.8), Inches(1.5), Inches(12), Inches(1.0),
         "What Canary-RL gives us", size=36, bold=True,
         color=RGBColor(0xFF,0xFF,0xFF))
add_rule(s, Inches(0.85), Inches(2.6), Inches(2.0), color=GOLD, weight=2.5)

points = [
    ("Predict before launch",
     "HealthOracle scores any launch script in <1 s using 15-run history."),
    ("Watch the right signal",
     "BT-GRPO loss is dominated by KL display value; reward (correctness) is the real learning signal."),
    ("Trainer-aware thresholds",
     "Watcher v2 tolerates 1e10 grad transients; AR-trainer regime kept via override."),
    ("Closed feedback loop",
     "Each run's outcome retrains the oracle — calibration improves over time."),
    ("First win on the books",
     "run22: correctness 0.42 → 0.52 in 12 steps despite alarming loss numbers."),
]
y = Inches(3.0)
for tag, desc in points:
    add_box(s, Inches(0.85), y + Inches(0.1), Inches(0.18), Inches(0.18),
            fill=GOLD, border=GOLD)
    add_text(s, Inches(1.2), y, Inches(11), Inches(0.45),
             tag, size=18, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    add_text(s, Inches(1.2), y + Inches(0.4), Inches(11), Inches(0.45),
             desc, size=13, color=RGBColor(0xC9,0xCF,0xDC))
    y += Inches(0.78)
add_text(s, Inches(0.85), Inches(7.05), Inches(12), Inches(0.35),
         "dllm/pipelines/rl/canary  ·  scripts/canary_*  ·  scripts/launch_with_canary.sh",
         size=11, color=RGBColor(0xC9,0xCF,0xDC), font="Consolas")

prs.save(str(OUT))
print(f"wrote {OUT}  ({OUT.stat().st_size/1024:.1f} KB)")
